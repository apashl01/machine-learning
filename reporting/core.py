"""
Core PowerPoint generation functionality.

Provides DesignReviewGenerator class for creating design review presentations.
"""

from pathlib import Path
from datetime import datetime
from typing import Optional, List, Dict, Any
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# Standard slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color scheme
COLORS = {
    'primary': RGBColor(0, 51, 102),      # Dark blue
    'secondary': RGBColor(0, 102, 153),   # Medium blue
    'accent': RGBColor(255, 153, 0),      # Orange
    'success': RGBColor(0, 153, 76),      # Green
    'warning': RGBColor(255, 193, 7),     # Yellow
    'danger': RGBColor(204, 0, 0),        # Red
    'text': RGBColor(51, 51, 51),         # Dark gray
    'light': RGBColor(245, 245, 245),     # Light gray
}


@dataclass
class SlideContent:
    """Container for slide content."""
    title: str
    subtitle: Optional[str] = None
    bullets: List[str] = field(default_factory=list)
    image_path: Optional[str] = None
    table_data: Optional[List[List[str]]] = None
    notes: Optional[str] = None


class DesignReviewGenerator:
    """
    Generate PowerPoint design review presentations.

    Usage:
        gen = DesignReviewGenerator("EW System Design Review")
        gen.add_title_slide(subtitle="Preliminary Design Review", date="2025-01-15")
        gen.add_section_slide("RF Chain Analysis")
        gen.add_content_slide("Link Budget", bullets=[...], image_path="plot.png")
        gen.save("design_review.pptx")

        # With template:
        gen = DesignReviewGenerator("EW System Design Review", template="template.pptx")
    """

    def __init__(self, title: str, author: str = "Analysis Team",
                 template: Optional[str] = None):
        """
        Initialize the design review generator.

        Args:
            title: Presentation title
            author: Author/team name
            template: Optional path to PowerPoint template (.pptx or .potx)
        """
        self.title = title
        self.author = author
        self.template = template

        # Load template if provided, otherwise create blank presentation
        if template and Path(template).exists():
            self.prs = Presentation(template)
            print(f"Using PowerPoint template: {template}")
        else:
            self.prs = Presentation()
            # Set slide dimensions to widescreen (only for blank presentations)
            self.prs.slide_width = SLIDE_WIDTH
            self.prs.slide_height = SLIDE_HEIGHT

        self.slide_count = 0

    def add_title_slide(self, subtitle: Optional[str] = None,
                        date: Optional[str] = None):
        """Add a title slide."""
        if date is None:
            date = datetime.now().strftime("%B %d, %Y")

        # Use blank layout and build custom
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Add background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['primary']
        shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.8), Inches(12), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER

        # Date and author
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(12), Inches(1)
        )
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{date}\n{self.author}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER

        self.slide_count += 1
        return slide

    def add_section_slide(self, section_title: str,
                          section_subtitle: Optional[str] = None):
        """Add a section divider slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Accent bar on left
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SLIDE_HEIGHT
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['accent']
        shape.line.fill.background()

        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']

        if section_subtitle:
            p = tf.add_paragraph()
            p.text = section_subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS['text']

        self.slide_count += 1
        return slide

    def add_content_slide(self, title: str,
                          bullets: Optional[List[str]] = None,
                          image_path: Optional[str] = None,
                          image_width: float = 6.0,
                          notes: Optional[str] = None):
        """
        Add a content slide with optional bullets and image.

        Args:
            title: Slide title
            bullets: List of bullet points
            image_path: Path to image file
            image_width: Width of image in inches
            notes: Speaker notes
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLORS['primary']
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Content area
        content_top = Inches(1.3)

        if bullets and image_path:
            # Split layout: bullets on left, image on right
            self._add_bullets(slide, bullets, Inches(0.5), content_top,
                            Inches(5.5), Inches(5.5))
            self._add_image(slide, image_path, Inches(6.5), content_top,
                          width=Inches(image_width))
        elif bullets:
            # Full width bullets
            self._add_bullets(slide, bullets, Inches(0.5), content_top,
                            Inches(12), Inches(5.5))
        elif image_path:
            # Centered image
            self._add_image(slide, image_path, Inches(1.5), content_top,
                          width=Inches(10))

        # Notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self.slide_count += 1
        return slide

    def add_two_image_slide(self, title: str,
                            image1_path: str, image1_caption: str,
                            image2_path: str, image2_caption: str,
                            notes: Optional[str] = None):
        """Add a slide with two images side by side."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLORS['primary']
        title_bar.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Left image
        self._add_image(slide, image1_path, Inches(0.3), Inches(1.3),
                       width=Inches(6.2))
        cap1 = slide.shapes.add_textbox(
            Inches(0.3), Inches(6.3), Inches(6.2), Inches(0.5)
        )
        tf = cap1.text_frame
        p = tf.paragraphs[0]
        p.text = image1_caption
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

        # Right image
        self._add_image(slide, image2_path, Inches(6.8), Inches(1.3),
                       width=Inches(6.2))
        cap2 = slide.shapes.add_textbox(
            Inches(6.8), Inches(6.3), Inches(6.2), Inches(0.5)
        )
        tf = cap2.text_frame
        p = tf.paragraphs[0]
        p.text = image2_caption
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self.slide_count += 1
        return slide

    def add_table_slide(self, title: str,
                        headers: List[str],
                        rows: List[List[str]],
                        notes: Optional[str] = None):
        """Add a slide with a table."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLORS['primary']
        title_bar.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Create table
        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)

        table_width = Inches(12)
        table_height = Inches(0.4 * n_rows)

        table = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.5), Inches(1.5),
            table_width, table_height
        ).table

        # Set column widths
        col_width = table_width / n_cols
        for col in table.columns:
            col.width = int(col_width)

        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['secondary']
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER

                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['light']

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self.slide_count += 1
        return slide

    def add_metrics_slide(self, title: str,
                          metrics: Dict[str, str],
                          subtitle: Optional[str] = None):
        """Add a slide with key metrics in large boxes."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLORS['primary']
        title_bar.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Metrics boxes
        n_metrics = len(metrics)
        if n_metrics <= 3:
            cols = n_metrics
            rows = 1
        elif n_metrics <= 6:
            cols = 3
            rows = 2
        else:
            cols = 4
            rows = (n_metrics + 3) // 4

        box_width = Inches(12 / cols - 0.2)
        box_height = Inches(2)

        metric_items = list(metrics.items())
        for idx, (label, value) in enumerate(metric_items):
            row = idx // cols
            col = idx % cols

            left = Inches(0.5 + col * (12 / cols))
            top = Inches(1.5 + row * 2.5)

            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = COLORS['light']
            box.line.color.rgb = COLORS['secondary']

            # Value
            val_box = slide.shapes.add_textbox(
                left, top + Inches(0.3), box_width, Inches(1)
            )
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.alignment = PP_ALIGN.CENTER

            # Label
            lbl_box = slide.shapes.add_textbox(
                left, top + Inches(1.3), box_width, Inches(0.5)
            )
            tf = lbl_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text']
            p.alignment = PP_ALIGN.CENTER

        self.slide_count += 1
        return slide

    def _add_bullets(self, slide, bullets: List[str],
                     left, top, width, height):
        """Add bullet points to a slide with proper formatting.

        Task 7.1 Fix: Uses PowerPoint's native bullet formatting instead of
        manually prepending bullet characters. Supports sub-bullets via tab
        or leading spaces.
        """
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Detect sub-bullets (leading tab or 3+ spaces)
            if bullet.startswith('\t') or bullet.startswith('   '):
                p.level = 1
                p.text = bullet.strip()
            else:
                p.level = 0
                p.text = bullet

            p.font.size = Pt(18 if p.level == 0 else 16)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(12 if p.level == 0 else 8)

    def _add_image(self, slide, image_path: str,
                   left, top, width=None, height=None):
        """Add an image to a slide with auto-scaling.

        Task 7.2 Fix: Implements scale-to-fit logic to prevent plots from
        exceeding slide boundaries. Images are scaled proportionally to fit
        within max_width and max_height bounds.
        """
        path = Path(image_path)
        if not path.exists():
            # Add placeholder text if image not found
            box = slide.shapes.add_textbox(left, top, Inches(4), Inches(2))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"[Image not found: {path.name}]"
            p.font.size = Pt(14)
            p.font.italic = True
            return

        # Define max bounds for scaling (Task 7.2)
        max_width = Inches(10)
        max_height = Inches(5.5)

        if width:
            # Add with specified width
            pic = slide.shapes.add_picture(str(path), left, top, width=width)

            # Check if height exceeds max and rescale proportionally
            if pic.height > max_height:
                scale = max_height / pic.height
                pic.width = int(pic.width * scale)
                pic.height = max_height
        elif height:
            pic = slide.shapes.add_picture(str(path), left, top, height=height)

            # Check if width exceeds max and rescale proportionally
            if pic.width > max_width:
                scale = max_width / pic.width
                pic.height = int(pic.height * scale)
                pic.width = max_width
        else:
            pic = slide.shapes.add_picture(str(path), left, top)

            # Scale to fit within bounds
            if pic.width > max_width or pic.height > max_height:
                width_scale = max_width / pic.width if pic.width > max_width else 1.0
                height_scale = max_height / pic.height if pic.height > max_height else 1.0
                scale = min(width_scale, height_scale)
                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)

    def save(self, output_path: str):
        """Save the presentation to a file."""
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output))
        print(f"Presentation saved: {output}")
        print(f"  Total slides: {self.slide_count}")
        return output
